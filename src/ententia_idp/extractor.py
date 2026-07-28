import io
from pathlib import Path
from typing import Any

import fitz
from docx import Document as DocxDocument
from pptx import Presentation

from .logger import logger
from .models import DocumentMetadata, ExtractedDocument, ExtractedFigure, ExtractedPage, ExtractedTable


def _format_table_markdown(rows: list[list[str]]) -> str:
    if not rows:
        return ""
    header = rows[0]
    lines = ["| " + " | ".join(cell or "" for cell in header) + " |"]
    lines.append("| " + " | ".join(["---"] * len(header)) + " |")
    for row in rows[1:]:
        lines.append("| " + " | ".join(cell or "" for cell in row) + " |")
    return "\n".join(lines)


def _format_table_csv(rows: list[list[str]]) -> str:
    if not rows:
        return ""
    lines: list[str] = []
    for row in rows:
        quoted = [f'"{cell.replace(chr(34), chr(34) * 2)}"' for cell in row]
        lines.append(",".join(quoted))
    return "\n".join(lines)


def _extract_pdf(file_path: str) -> ExtractedDocument:
    logger.info("Running PDF extraction for %s", file_path)
    doc = fitz.open(file_path)
    pages: list[ExtractedPage] = []
    figures: list[ExtractedFigure] = []

    for page_number, page in enumerate(doc, start=1):
        text = page.get_text("text").strip()
        image_count = 0

        for img_info in page.get_images(full=True):
            xref = img_info[0]
            try:
                pix = fitz.Pixmap(doc, xref)
                if pix.n > 4:
                    pix = fitz.Pixmap(fitz.csRGB, pix)
                png_bytes = pix.tobytes("png")
            except Exception:
                continue

            if len(png_bytes) == 0:
                continue

            image_count += 1
            figures.append(ExtractedFigure(page=page_number, caption="", fig_id=f"pdf-{page_number}-{xref}"))

        pages.append(
            ExtractedPage(
                page_number=page_number,
                text=text,
                word_count=len(text.split()),
                figure_count=image_count,
                table_count=0,
            )
        )

    metadata = DocumentMetadata(
        source_file_name=Path(file_path).name,
        file_type=Path(file_path).suffix.lstrip(".").upper(),
        total_pages=len(doc),
    )

    return ExtractedDocument(metadata=metadata, pages=pages, tables=[], figures=figures)


def _extract_docx(file_path: str) -> ExtractedDocument:
    logger.info("Running DOCX extraction for %s", file_path)
    document = DocxDocument(file_path)
    text = "\n\n".join(paragraph.text for paragraph in document.paragraphs if paragraph.text)

    tables: list[ExtractedTable] = []
    for table in document.tables:
        rows = [[cell.text or "" for cell in row.cells] for row in table.rows]
        tables.append(
            ExtractedTable(
                page=1,
                markdown=_format_table_markdown(rows),
                csv=_format_table_csv(rows),
            )
        )

    figures: list[ExtractedFigure] = []
    image_index = 0
    for rel in document.part._rels.values():
        if "image" in rel.reltype:
            image_index += 1
            figures.append(ExtractedFigure(page=1, caption="", fig_id=f"docx-{image_index}"))

    pages = [
        ExtractedPage(
            page_number=1,
            text=text,
            word_count=len(text.split()),
            figure_count=len(figures),
            table_count=len(tables),
        )
    ]

    metadata = DocumentMetadata(
        source_file_name=Path(file_path).name,
        file_type=Path(file_path).suffix.lstrip(".").upper(),
        total_pages=1,
    )

    return ExtractedDocument(metadata=metadata, pages=pages, tables=tables, figures=figures)


def _extract_pptx(file_path: str) -> ExtractedDocument:
    logger.info("Running PPTX extraction for %s", file_path)
    presentation = Presentation(file_path)
    pages: list[ExtractedPage] = []
    figures: list[ExtractedFigure] = []

    for slide_number, slide in enumerate(presentation.slides, start=1):
        slide_texts: list[str] = []
        image_count = 0

        for shape in slide.shapes:
            if shape.has_text_frame:
                slide_texts.append(shape.text)
            if shape.shape_type == 13:  # picture
                image_count += 1
                figures.append(ExtractedFigure(page=slide_number, caption="", fig_id=f"pptx-{slide_number}-{image_count}"))

        page_text = "\n\n".join([text for text in slide_texts if text])
        pages.append(
            ExtractedPage(
                page_number=slide_number,
                text=page_text,
                word_count=len(page_text.split()),
                figure_count=image_count,
                table_count=0,
            )
        )

    metadata = DocumentMetadata(
        source_file_name=Path(file_path).name,
        file_type=Path(file_path).suffix.lstrip(".").upper(),
        total_pages=len(presentation.slides),
    )

    return ExtractedDocument(metadata=metadata, pages=pages, tables=[], figures=figures)


def _extract_txt(file_path: str) -> ExtractedDocument:
    logger.info("Running plain text extraction for %s", file_path)
    with open(file_path, "r", encoding="utf-8", errors="replace") as handle:
        text = handle.read().strip()

    pages = [
        ExtractedPage(
            page_number=1,
            text=text,
            word_count=len(text.split()),
            figure_count=0,
            table_count=0,
        )
    ]

    metadata = DocumentMetadata(
        source_file_name=Path(file_path).name,
        file_type=Path(file_path).suffix.lstrip(".").upper(),
        total_pages=1,
    )

    return ExtractedDocument(metadata=metadata, pages=pages, tables=[], figures=[])


class DocumentExtractor:
    def extract(self, file_path: str) -> ExtractedDocument:
        suffix = Path(file_path).suffix.lower()
        if suffix == ".pdf":
            return _extract_pdf(file_path)
        if suffix == ".docx":
            return _extract_docx(file_path)
        if suffix == ".pptx":
            return _extract_pptx(file_path)
        if suffix == ".txt":
            return _extract_txt(file_path)

        raise ValueError(f"Unsupported file type: {suffix}")
